"""
AX Hackathon — Phase 1 Blueprint Builder (CONTENT-FILLED)
=========================================================

Renders the full 12-slide PS04 blueprint for submission.

Install:
    pip install python-pptx pillow

Run from the BlueprintPPT/ directory:
    python build_blueprint.py
    # -> AX_Hackathon_Phase1_Blueprint_FILLED.pptx
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ──────────────────────────────────────────────────────────────────────────────
# DESIGN TOKENS
# ──────────────────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x00, 0x3B, 0x8E)
SRI_BLUE    = RGBColor(0x14, 0x4B, 0xC4)
SRI_DARK    = RGBColor(0x0F, 0x1B, 0x3D)
SRI_RED     = RGBColor(0xC3, 0x29, 0x2B)
SRI_GOLD    = RGBColor(0xB8, 0x86, 0x0B)
SRI_GREEN   = RGBColor(0x1B, 0x7F, 0x4F)
SRI_GRAY    = RGBColor(0x6A, 0x73, 0x83)
SLIDE_BG    = RGBColor(0xFA, 0xFC, 0xFF)
CARD_FILL   = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xD9, 0xDE, 0xE6)
LABEL_BLUE  = NAVY
BODY_GRAY   = RGBColor(0x40, 0x48, 0x56)
KPI_COLOR   = SRI_BLUE
NOVEL_COLOR = SRI_RED
ACCENT_BG_BLUE  = RGBColor(0xEC, 0xF2, 0xFC)
ACCENT_BG_RED   = RGBColor(0xFD, 0xEE, 0xEE)
ACCENT_BG_GOLD  = RGBColor(0xFB, 0xF2, 0xDD)
ACCENT_BG_GREEN = RGBColor(0xE7, 0xF4, 0xEC)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.85)
MARGIN_X = Inches(0.4)

ILLUS_DIR = Path(__file__).resolve().parent.parent / "illustrations"


# ──────────────────────────────────────────────────────────────────────────────
# LOW-LEVEL HELPERS
# ──────────────────────────────────────────────────────────────────────────────
def set_bg(slide, rgb):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = rgb


def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    return s


def add_card(slide, x, y, w, h, fill=CARD_FILL, border=CARD_BORDER, radius=0.06):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = border
    c.line.width = Pt(0.75)
    c.adjustments[0] = radius
    c.text_frame.text = ""
    return c


def add_text(slide, x, y, w, h, paragraphs, anchor="top", margin=0.08):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin * 0.6)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[anchor]
    for i, p in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = {
            "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
        }[p.get("align", "left")]
        if "space_after" in p:
            para.space_after = Pt(p["space_after"])
        if "space_before" in p:
            para.space_before = Pt(p["space_before"])
        if "level" in p:
            para.level = p["level"]
        # Each paragraph can be a list of runs (rich formatting) or a single "text"
        runs = p.get("runs")
        if runs is None:
            runs = [{
                "text": p.get("text", ""),
                "size": p.get("size", 12),
                "bold": p.get("bold", False),
                "italic": p.get("italic", False),
                "color": p.get("color", BODY_GRAY),
            }]
        for r in runs:
            run = para.add_run()
            run.text = r.get("text", "")
            run.font.name = FONT
            run.font.size = Pt(r.get("size", 12))
            run.font.bold = r.get("bold", False)
            run.font.italic = r.get("italic", False)
            run.font.color.rgb = r.get("color", BODY_GRAY)
    return tb


def add_header(slide, title):
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, NAVY)
    add_text(
        slide, Inches(0.4), 0, SLIDE_W - Inches(0.8), HEADER_H,
        [{"text": title, "size": 24, "bold": True,
          "color": RGBColor(0xFF, 0xFF, 0xFF)}],
        anchor="middle",
    )
    # thin gold accent under header
    add_rect(slide, 0, HEADER_H, SLIDE_W, Inches(0.04), SRI_GOLD)


def add_footer(slide, page_num, total=18):
    add_text(
        slide,
        SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.35),
        Inches(1.3), Inches(0.3),
        [{"text": f"{page_num} / {total}", "size": 9,
          "color": SRI_GRAY, "align": "right"}],
    )


def add_image_fit(slide, path, x, y, max_w, max_h):
    """Place image fitted inside (max_w, max_h) keeping aspect ratio, centered."""
    im = Image.open(path)
    iw, ih = im.size
    ar = iw / ih
    # Convert max_w, max_h (EMU) to inches for math, then back
    max_w_in = max_w / 914400
    max_h_in = max_h / 914400
    w_in = max_w_in
    h_in = w_in / ar
    if h_in > max_h_in:
        h_in = max_h_in
        w_in = h_in * ar
    cx = x + (max_w - Inches(w_in)) / 2
    cy = y + (max_h - Inches(h_in)) / 2
    slide.shapes.add_picture(str(path), cx, cy, Inches(w_in), Inches(h_in))


def blank_slide(prs, title, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SLIDE_BG)
    add_header(slide, title)
    add_footer(slide, page)
    return slide


def kpi_run(text, size=14, color=KPI_COLOR):
    return {"text": text, "size": size, "bold": True, "color": color}


def label_para(text, size=14):
    return {"text": text, "size": size, "bold": True, "color": LABEL_BLUE,
            "space_after": 4}


def body_para(text, size=11, color=BODY_GRAY, space_after=2, **kw):
    return {"text": text, "size": size, "color": color,
            "space_after": space_after, **kw}


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ──────────────────────────────────────────────────────────────────────────────
def slide_1_team(prs, page):
    s = blank_slide(prs, "Team Introduction", page)

    # Team name card (full width)
    add_text(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
             [label_para("Team Name", 16)])
    add_card(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.9))
    add_text(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(0.9),
             [{"text": "Team MalayM09 — sole-team submission",
               "size": 18, "bold": True, "color": SRI_DARK}],
             anchor="middle")

    # Member card (full width since only 1 member)
    add_text(s, Inches(0.5), Inches(2.65), Inches(12.3), Inches(0.4),
             [label_para("Team Member", 16)])
    add_card(s, Inches(0.5), Inches(3.10), Inches(12.3), Inches(3.8))
    paras = [
        {"text": "Malay Mishra", "size": 22, "bold": True,
         "color": SRI_DARK, "space_after": 14},
        {"runs": [
            {"text": "College  ", "size": 13, "bold": True, "color": LABEL_BLUE},
            {"text": "BITS Pilani", "size": 13, "color": BODY_GRAY},
        ], "space_after": 8},
        {"runs": [
            {"text": "Department  ", "size": 13, "bold": True, "color": LABEL_BLUE},
            {"text": "Computer Science",
             "size": 13, "color": BODY_GRAY},
        ], "space_after": 8},
        {"runs": [
            {"text": "Year  ", "size": 13, "bold": True, "color": LABEL_BLUE},
            {"text": "Final year", "size": 13, "color": BODY_GRAY},
        ], "space_after": 8},
        {"runs": [
            {"text": "Email  ", "size": 13, "bold": True, "color": LABEL_BLUE},
            {"text": "f20220116@pilani.bits-pilani.ac.in", "size": 13, "color": BODY_GRAY},
        ], "space_after": 8},
    ]
    add_text(s, Inches(0.9), Inches(3.30), Inches(11.5), Inches(3.4), paras)


def slide_2_problem(prs, page):
    s = blank_slide(prs, "Problem Statement", page)

    # Selected problem (short card)
    add_text(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
             [label_para("Selected Problem", 15)])
    add_card(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.7),
             fill=ACCENT_BG_BLUE, border=SRI_BLUE)
    add_text(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(0.7),
             [{"text": "PS04  ·  Speech Disentanglement — "
                       "Speaker-Specific Custom Word Detection on the Edge",
               "size": 15, "bold": True, "color": SRI_DARK}],
             anchor="middle")

    # Problem Understanding
    add_text(s, Inches(0.5), Inches(2.30), Inches(12.3), Inches(0.4),
             [label_para("Problem Understanding", 15)])
    add_card(s, Inches(0.5), Inches(2.75), Inches(12.3), Inches(2.2))
    body = [
        body_para("Voice triggers must wake a device only when the right "
                  "person says the right word — not when an impostor "
                  "says it, not when a phonetically similar word is spoken, "
                  "and not when crowd / traffic / babble noise drowns the "
                  "signal.", size=12, space_after=6),
        body_para("Today's edge wake-word systems either (a) accept anyone "
                  "saying the trigger (privacy risk on shared devices), or "
                  "(b) need >10 M parameters and cloud offload. Standard "
                  "log-Mel + ResNet pipelines collapse at low SNR; "
                  "Transformer-based KWS miss the latency budget.",
                  size=12, space_after=6),
    ]
    add_text(s, Inches(0.7), Inches(2.85), Inches(12.0), Inches(2.0), body)

    # KPI strip (3 columns)
    add_text(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.4),
             [label_para("Hard KPI envelope", 15)])
    kpi_y = Inches(5.5)
    kpi_h = Inches(1.6)
    kpi_w = Inches(3.95)
    kpi_gap = Inches(0.2)
    kpi_x = [Inches(0.5),
             Inches(0.5) + kpi_w + kpi_gap,
             Inches(0.5) + (kpi_w + kpi_gap) * 2]
    titles = [
        ("Accuracy", [
            ("True-accept clean", "≥ 99 %"),
            ("True-accept noisy", "≥ 90 %"),
            ("False-accept rate", "< 1 / hr"),
        ]),
        ("Compute", [
            ("Params", "< 3 M"),
            ("xRT (real-time)", "< 0.2 s"),
            ("Memory", "edge-class"),
        ]),
        ("Conditions", [
            ("SNR sweep", "−5 to 30 dB"),
            ("Distance", "0.5 to 5 m"),
            ("Languages", "phonetic-general"),
        ]),
    ]
    for x, (head, rows) in zip(kpi_x, titles):
        add_card(s, x, kpi_y, kpi_w, kpi_h,
                 fill=ACCENT_BG_BLUE, border=SRI_BLUE)
        paras = [{"text": head, "size": 13, "bold": True, "color": SRI_BLUE,
                  "space_after": 6}]
        for lbl, val in rows:
            paras.append({"runs": [
                {"text": f"{lbl}   ", "size": 11, "color": BODY_GRAY},
                {"text": val, "size": 12, "bold": True, "color": SRI_DARK},
            ], "space_after": 3})
        add_text(s, x + Inches(0.2), kpi_y + Inches(0.12),
                 kpi_w - Inches(0.4), kpi_h - Inches(0.2), paras)


def slide_3_solution_high(prs, page):
    s = blank_slide(prs, "Proposed Solution — High Level", page)

    # One-sentence summary
    add_card(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.95),
             fill=ACCENT_BG_BLUE, border=SRI_BLUE)
    add_text(s, Inches(0.7), Inches(1.05), Inches(12.0), Inches(0.95),
             [{"runs": [
                 {"text": "One sentence.  ", "size": 13, "bold": True,
                  "color": SRI_BLUE},
                 {"text": "A single ~830 K-parameter joint student "
                          "(KWS backbone validated at 324 K) performs Keyword "
                          "Spotting and Speaker Verification, distilled from "
                          "two frozen SOTA teachers, with speaker-conditioned "
                          "classification and a FAR-aware training loop.",
                  "size": 13, "color": SRI_DARK},
             ]}],
             anchor="middle")

    # 3 columns of pillars
    col_y = Inches(2.25)
    col_h = Inches(4.5)
    col_w = Inches(4.0)
    col_gap = Inches(0.13)
    col_x = [Inches(0.5),
             Inches(0.5) + col_w + col_gap,
             Inches(0.5) + (col_w + col_gap) * 2]

    pillars = [
        ("1. Robust Frontend", SRI_BLUE,
         "Learnable PCEN on Mel — survives −5 dB SNR where log-Mel "
         "collapses. Equivalent to LEAF in practice at a fraction of the "
         "compute.",
         "PCEN ≡ LEAF (sub-layer)"),
        ("2. Compact Joint Student", SRI_GOLD,
         "Modified BC-ResNet-8 with TRM branching, MQMHA pooling on SV head, "
         "and FiLM-conditioned KWS head — the speaker embedding modulates "
         "the keyword classifier.",
         "499 K Phase-1, → ~830 K Phase-2"),
        ("3. Multi-Teacher KD + FAR Loop", SRI_RED,
         "WavLM-Base+ (phonetic) + ECAPA-TDNN (biometric) teachers. Training "
         "loop includes online τₛ calibration that bakes the "
         "FA < 1 / hr operating point into the embedding geometry.",
         "FA < 1 / hr baked into training"),
    ]
    for x, (title, color, body, tag) in zip(col_x, pillars):
        add_card(s, x, col_y, col_w, col_h)
        # colored top stripe inside the card
        add_rect(s, x, col_y, col_w, Inches(0.35), color)
        add_text(s, x + Inches(0.2), col_y, col_w - Inches(0.4), Inches(0.35),
                 [{"text": title, "size": 14, "bold": True,
                   "color": RGBColor(0xFF, 0xFF, 0xFF), "align": "left"}],
                 anchor="middle")
        # body
        add_text(s, x + Inches(0.25), col_y + Inches(0.5),
                 col_w - Inches(0.5), col_h - Inches(1.1),
                 [{"text": body, "size": 12, "color": BODY_GRAY}])
        # tag at bottom
        add_text(s, x + Inches(0.25), col_y + col_h - Inches(0.55),
                 col_w - Inches(0.5), Inches(0.4),
                 [{"text": tag, "size": 11, "bold": True, "italic": True,
                   "color": color}])


def slide_4_architecture(prs, page):
    s = blank_slide(prs, "Technical Details — System Architecture", page)

    # Image (wide aspect 5.94)
    img_x = Inches(0.4)
    img_y = Inches(1.05)
    img_max_w = Inches(12.5)
    img_max_h = Inches(2.6)
    add_image_fit(s, ILLUS_DIR / "1_architecture.png",
                  img_x, img_y, img_max_w, img_max_h)

    # Three annotation cards below (Forward / Training / Params)
    ann_y = Inches(3.85)
    ann_h = Inches(3.15)
    ann_w = Inches(4.0)
    ann_gap = Inches(0.13)
    ann_x = [Inches(0.5),
             Inches(0.5) + ann_w + ann_gap,
             Inches(0.5) + (ann_w + ann_gap) * 2]

    annotations = [
        ("Forward (inference)", SRI_BLUE, ACCENT_BG_BLUE, [
            "waveform → PCEN frontend",
            "→ shared BC-ResNet-8 trunk (blue)",
            "→ TRM branching (yellow)",
            "→ KWS / SV heads (green)",
            "→ FiLM (red) modulates KWS logits on the speaker embedding.",
        ]),
        ("Training", SRI_RED, ACCENT_BG_RED, [
            "Stop-gradient on the SV→FiLM edge —",
            "KWS loss never back-propagates into the speaker branch.",
            "Otherwise SV degrades (trade-off measured on Slide 7).",
        ]),
        ("Params", SRI_GOLD, ACCENT_BG_GOLD, [
            "Trunk  324 K  (validated)",
            "+ branches / heads  ~175 K",
            "=  499 K  Phase-1 joint  (Slide 8)",
            "~830 K projected with multi-teacher KD heads.",
        ]),
    ]
    for x, (title, color, bg, lines) in zip(ann_x, annotations):
        add_card(s, x, ann_y, ann_w, ann_h, fill=bg, border=color)
        add_text(s, x + Inches(0.2), ann_y + Inches(0.1),
                 ann_w - Inches(0.4), Inches(0.4),
                 [{"text": title, "size": 13, "bold": True, "color": color}])
        body_paras = [{"text": ln, "size": 11, "color": BODY_GRAY,
                       "space_after": 3} for ln in lines]
        add_text(s, x + Inches(0.25), ann_y + Inches(0.55),
                 ann_w - Inches(0.5), ann_h - Inches(0.7), body_paras)


def slide_5_distillation(prs, page):
    s = blank_slide(prs, "Technical Details — Multi-Teacher Knowledge Distillation", page)

    # Image (aspect 2.14, wider than tall)
    img_x = Inches(0.4)
    img_y = Inches(1.05)
    img_max_w = Inches(12.5)
    img_max_h = Inches(3.6)
    add_image_fit(s, ILLUS_DIR / "2_distillation.png",
                  img_x, img_y, img_max_w, img_max_h)

    # Loss equation card
    add_card(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(1.0),
             fill=ACCENT_BG_BLUE, border=SRI_BLUE)
    add_text(s, Inches(0.7), Inches(4.85), Inches(12.0), Inches(1.0),
             [{"runs": [
                 {"text": "L_total = ", "size": 14, "bold": True,
                  "color": SRI_DARK},
                 {"text": "λ₁·L_KWS^Focal "
                          "+ λ₂·L_SV^SubCtr-AAM "
                          "+ λ₃·T²·KL(stu ∥ probe) "
                          "+ λ₄·L_hint^{4,8,11} "
                          "+ λ₅·L_ECAPA^multi-level",
                  "size": 13, "color": SRI_DARK},
             ]}],
             anchor="middle")

    # Phase-1 validation strip
    add_card(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.1),
             fill=ACCENT_BG_GREEN, border=SRI_GREEN)
    add_text(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(1.1),
             [{"runs": [
                 {"text": "Phase-1 validation:  ", "size": 13, "bold": True,
                  "color": SRI_GREEN},
                 {"text": "WavLM-Base+ final-layer distillation (simplest "
                          "form of L_hint) trained jointly with L_KWS for "
                          "40 epochs  →  KWS clean val ",
                  "size": 12, "color": BODY_GRAY},
                 {"text": "97.6 → 98.3 %", "size": 13, "bold": True,
                  "color": SRI_BLUE},
                 {"text": ",  cosine alignment with teacher ", "size": 12,
                  "color": BODY_GRAY},
                 {"text": "0.94", "size": 13, "bold": True, "color": SRI_BLUE},
                 {"text": ".  Multi-layer hints {4, 8, 11} projected to close "
                          "the remaining 0.7 pp gap to 99 % — see Slide 7.",
                  "size": 12, "color": BODY_GRAY},
             ]}],
             anchor="middle")


def slide_6_params_stages(prs, page):
    s = blank_slide(prs, "Technical Details — Parameter Budget & Training Stages", page)

    # Left: param budget figure
    img_max_w = Inches(6.0)
    img_max_h = Inches(5.6)
    add_image_fit(s, ILLUS_DIR / "6_param_budget.png",
                  Inches(0.4), Inches(1.1), img_max_w, img_max_h)

    # Right: training stages table
    tbl_x = Inches(6.7)
    tbl_y = Inches(1.1)
    tbl_w = Inches(6.3)
    tbl_h = Inches(5.6)
    add_card(s, tbl_x, tbl_y, tbl_w, tbl_h)

    add_text(s, tbl_x + Inches(0.2), tbl_y + Inches(0.1),
             tbl_w - Inches(0.4), Inches(0.45),
             [{"text": "Training stages  (~36 GPU-hr Kaggle T4 / P100)",
               "size": 13, "bold": True, "color": SRI_DARK}])

    stages = [
        ("0",  "SC + LibriPhrase",           "Train WavLM→KWS probe (one-time)"),
        ("1",  "Speech Commands V2",          "Stabilise PCEN + trunk"),
        ("2",  "SC + VoxCeleb1 + LibriPhrase","Full multi-task + all KD terms"),
        ("3",  "LibriPhrase + TTS-trigger",   "Specialise + online τₛ calibration"),
        ("4",  "Same",                        "Quantisation-aware training (INT8)"),
    ]
    row_h = Inches(0.55)
    row_y = tbl_y + Inches(0.7)
    # header row
    add_rect(s, tbl_x + Inches(0.2), row_y, tbl_w - Inches(0.4), Inches(0.4),
             SRI_BLUE)
    headers = [("Stage", Inches(0.6)),
               ("Data", Inches(2.6)),
               ("Purpose", Inches(2.9))]
    cx = tbl_x + Inches(0.3)
    for hh, hw in headers:
        add_text(s, cx, row_y, hw, Inches(0.4),
                 [{"text": hh, "size": 11, "bold": True,
                   "color": RGBColor(0xFF, 0xFF, 0xFF)}], anchor="middle")
        cx += hw
    row_y += Inches(0.42)
    for i, (stage, data, purp) in enumerate(stages):
        if i % 2 == 0:
            add_rect(s, tbl_x + Inches(0.2), row_y, tbl_w - Inches(0.4),
                     row_h, RGBColor(0xF4, 0xF7, 0xFB))
        cx = tbl_x + Inches(0.3)
        for val, ww, bold in [(stage, Inches(0.6), True),
                               (data, Inches(2.6), False),
                               (purp, Inches(2.9), False)]:
            add_text(s, cx, row_y, ww, row_h,
                     [{"text": val, "size": 11, "bold": bold,
                       "color": SRI_DARK if bold else BODY_GRAY}],
                     anchor="middle")
            cx += ww
        row_y += row_h

    # Bottom callout
    add_text(s, tbl_x + Inches(0.2), tbl_y + tbl_h - Inches(0.85),
             tbl_w - Inches(0.4), Inches(0.75),
             [{"runs": [
                 {"text": "Final inference model:  ", "size": 12,
                  "bold": True, "color": SRI_GREEN},
                 {"text": "499.6 K params (17 % of 3 M cap), "
                          "INT8 ONNX ships at 0.79 MB.",
                  "size": 12, "color": BODY_GRAY},
             ]}])


def slide_7_phase1_kws(prs, page):
    s = blank_slide(prs,
                    "Phase-1 Empirical — KWS Backbone + Custom-Word Generalization",
                    page)

    # Left: trajectory + figure 10 (custom word)
    left_w = Inches(7.4)
    add_card(s, Inches(0.4), Inches(1.05), left_w, Inches(2.1),
             fill=ACCENT_BG_GREEN, border=SRI_GREEN)
    add_text(s, Inches(0.6), Inches(1.05), left_w - Inches(0.4), Inches(2.1),
             [{"text": "Measured path to ≥99 % TA Clean", "size": 13,
               "bold": True, "color": SRI_GREEN, "space_after": 4},
              {"runs": [
                  {"text": "97.6 %", "size": 16, "bold": True,
                   "color": SRI_BLUE},
                  {"text": "  (P0, 322 K params)   →   ",
                   "size": 12, "color": BODY_GRAY},
                  {"text": "98.2 %", "size": 16, "bold": True,
                   "color": SRI_BLUE},
                  {"text": "  (P2 wide 705 K + label-smooth)   →   ",
                   "size": 12, "color": BODY_GRAY},
                  {"text": "98.3 %", "size": 16, "bold": True,
                   "color": SRI_BLUE},
                  {"text": "  (P3 + WavLM final-layer distill)   →   ",
                   "size": 12, "color": BODY_GRAY},
                  {"text": "≥99 %", "size": 16, "bold": True,
                   "color": SRI_GREEN},
                  {"text": "  projected (multi-layer SKILL hints + EMA).",
                   "size": 12, "color": BODY_GRAY},
              ], "space_after": 6},
              {"text": "All within < 24 % of the 3 M param cap.",
               "size": 11, "italic": True, "color": SRI_GRAY}],
             anchor="middle")

    # Custom-word figure below
    add_image_fit(s, ILLUS_DIR / "10_custom_word.png",
                  Inches(0.4), Inches(3.25), left_w, Inches(3.85))

    # Right: 3 stacked headline cards (English EER, Cross-lingual, Path to GREEN)
    rx = Inches(8.0)
    rw = Inches(4.95)

    # Card 1: Headline
    add_card(s, rx, Inches(1.05), rw, Inches(1.55),
             fill=ACCENT_BG_BLUE, border=SRI_BLUE)
    add_text(s, rx + Inches(0.2), Inches(1.1), rw - Inches(0.4), Inches(1.5),
             [{"text": "Headline", "size": 13, "bold": True,
               "color": SRI_BLUE, "space_after": 4},
              {"runs": [
                  {"text": "EER = 12.43 %", "size": 14, "bold": True,
                   "color": SRI_BLUE},
                  {"text": "  on 7 novel SC V2 keywords (700 trials).",
                   "size": 11, "color": BODY_GRAY},
              ], "space_after": 2},
              {"runs": [
                  {"text": "AUC = 0.934", "size": 12, "bold": True,
                   "color": SRI_BLUE},
                  {"text": ", score separation ", "size": 11, "color": BODY_GRAY},
                  {"text": "4.2×", "size": 12, "bold": True, "color": SRI_BLUE},
                  {"text": " — at 384 K params, no WavLM teacher.",
                   "size": 11, "color": BODY_GRAY},
              ]}])

    # Card 2: Phonetic + cross-lingual (H1 evidence)
    add_card(s, rx, Inches(2.75), rw, Inches(2.55),
             fill=ACCENT_BG_RED, border=SRI_RED)
    add_text(s, rx + Inches(0.2), Inches(2.8), rw - Inches(0.4), Inches(2.5),
             [{"text": "Phonetic + cross-lingual generalisation",
               "size": 13, "bold": True, "color": SRI_RED, "space_after": 4},
              {"runs": [
                  {"text": "Within English, per-word EER 6–18 %  ",
                   "size": 11, "color": BODY_GRAY},
                  {"text": "(forward / follow at 6–7 %; backward / learn at "
                           "16–18 %).",
                   "size": 11, "color": BODY_GRAY},
              ], "space_after": 4},
              {"runs": [
                  {"text": "Zero-shot Tamil ", "size": 11, "bold": True,
                   "color": BODY_GRAY},
                  {"text": "(MSWC, 6 minimal-pair keywords, n=180): ",
                   "size": 11, "color": BODY_GRAY},
                  {"text": "14.7 % mean EER", "size": 13, "bold": True,
                   "color": SRI_RED},
                  {"text": "  —  only ", "size": 11, "color": BODY_GRAY},
                  {"text": "+2.2 pp", "size": 13, "bold": True,
                   "color": SRI_RED},
                  {"text": " above English.", "size": 11, "color": BODY_GRAY},
              ], "space_after": 4},
              {"text": "The embedder is phonetically organised, not "
                       "language-specific — precondition for user-chosen "
                       "triggers in any language.",
               "size": 11, "italic": True, "color": BODY_GRAY}])

    # Card 3: Path to GREEN
    add_card(s, rx, Inches(5.4), rw, Inches(1.75),
             fill=ACCENT_BG_GOLD, border=SRI_GOLD)
    add_text(s, rx + Inches(0.2), Inches(5.45), rw - Inches(0.4), Inches(1.7),
             [{"text": "Path to GREEN  (≤10 % EER)", "size": 13,
               "bold": True, "color": SRI_GOLD, "space_after": 4},
              {"text": "Add WavLM hint-layer distillation (layers 4 / 8 / 11). "
                       "Projected: ≤8 % EER, TAR ≥ 70 % at FAR = 1 %. "
                       "Same 384 K student, no extra inference cost.",
               "size": 11, "color": BODY_GRAY}])


def slide_8_phase1_deploy_babble(prs, page):
    s = blank_slide(prs,
                    "Phase-1 Empirical — Deployment Latency & Cocktail-Party",
                    page)

    # Left: deployment figure + side card
    add_image_fit(s, ILLUS_DIR / "11_deploy_metrics.png",
                  Inches(0.4), Inches(1.05), Inches(7.3), Inches(3.0))

    add_card(s, Inches(0.4), Inches(4.2), Inches(7.3), Inches(2.95),
             fill=ACCENT_BG_BLUE, border=SRI_BLUE)
    add_text(s, Inches(0.6), Inches(4.25), Inches(7.0), Inches(2.9),
             [{"text": "xRT  ·  comfortable PASS", "size": 13, "bold": True,
               "color": SRI_BLUE, "space_after": 4},
              {"runs": [
                  {"text": "Pi 4B INT8: ", "size": 11, "color": BODY_GRAY},
                  {"text": "12 ms", "size": 13, "bold": True, "color": SRI_BLUE},
                  {"text": " end-to-end vs ", "size": 11, "color": BODY_GRAY},
                  {"text": "200 ms", "size": 13, "bold": True, "color": SRI_BLUE},
                  {"text": " budget — ", "size": 11, "color": BODY_GRAY},
                  {"text": "17× headroom", "size": 13, "bold": True,
                   "color": SRI_GREEN},
                  {"text": ". Even FP32 at 17 ms clears 12×.",
                   "size": 11, "color": BODY_GRAY},
              ], "space_after": 6},
              {"text": "Param + size budget",
               "size": 12, "bold": True, "color": SRI_DARK, "space_after": 3},
              {"runs": [
                  {"text": "499.6 K params (17 % of 3 M cap), INT8 ONNX ships "
                           "at ", "size": 11, "color": BODY_GRAY},
                  {"text": "0.79 MB", "size": 12, "bold": True,
                   "color": SRI_BLUE},
                  {"text": " (2.6× shrink from 2.05 MB FP32).",
                   "size": 11, "color": BODY_GRAY},
              ], "space_after": 6},
              {"text": "Numeric fidelity",
               "size": 12, "bold": True, "color": SRI_DARK, "space_after": 3},
              {"text": "FP32 ONNX vs PyTorch max-abs-diff < 3×10⁻⁶ "
                       "(essentially exact); INT8 drift compensated by QAT "
                       "(Stage 4).",
               "size": 11, "color": BODY_GRAY}])

    # Right: babble figure + side card
    add_image_fit(s, ILLUS_DIR / "12_babble_robustness.png",
                  Inches(7.95), Inches(1.05), Inches(5.0), Inches(3.0))

    add_card(s, Inches(7.95), Inches(4.2), Inches(5.0), Inches(2.95),
             fill=ACCENT_BG_GREEN, border=SRI_GREEN)
    add_text(s, Inches(8.15), Inches(4.25), Inches(4.7), Inches(2.9),
             [{"text": "Cocktail-party robustness", "size": 13, "bold": True,
               "color": SRI_GREEN, "space_after": 4},
              {"runs": [
                  {"text": "Babble augmentation (40 % of batches, 1–2 voices) "
                           "lifts ", "size": 11, "color": BODY_GRAY},
                  {"text": "+19 pp at −5 dB", "size": 12, "bold": True,
                   "color": SRI_GREEN},
                  {"text": ", ", "size": 11, "color": BODY_GRAY},
                  {"text": "+21 pp at 0 dB", "size": 12, "bold": True,
                   "color": SRI_GREEN},
                  {"text": "; clean cost only −0.5 pp.", "size": 11,
                   "color": BODY_GRAY},
              ], "space_after": 6},
              {"runs": [
                  {"text": "N = 3 voices generalisation holds: ",
                   "size": 11, "color": BODY_GRAY},
                  {"text": "90 % at 10 dB", "size": 13, "bold": True,
                   "color": SRI_GREEN},
                  {"text": " with three simultaneous speakers — within KPI "
                           "envelope across the realistic-noise band.",
                   "size": 11, "color": BODY_GRAY},
              ]}])


def slide_9_novelty(prs, page):
    s = blank_slide(prs, "Novelty & Innovation", page)

    # Top card: core novelty
    add_card(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(1.95),
             fill=ACCENT_BG_RED, border=SRI_RED)
    add_text(s, Inches(0.7), Inches(1.1), Inches(12.0), Inches(1.9),
             [{"text": "Core novelty  ·  structural false-accept rejection via "
                       "FiLM-conditioned KWS",
               "size": 14, "bold": True, "color": SRI_RED, "space_after": 6},
              {"text": "Conventional personalised wake-word systems run KWS "
                       "and SV as independent gates and AND the decisions — "
                       "any leakage between heads or threshold drift breaks "
                       "the FA < 1 / hr budget.",
               "size": 11, "color": BODY_GRAY, "space_after": 4},
              {"text": "We instead modulate the KWS classifier's intermediate "
                       "features with the enrolled user's voice embedding, so "
                       "the same audio produces different logits for "
                       "different enrolled users. Impostor rejection becomes "
                       "a property of the network, not of a tunable threshold.",
               "size": 11, "color": BODY_GRAY}])

    # 4 supporting novelty cards (2x2)
    sup_x = [Inches(0.5), Inches(6.7)]
    sup_y = [Inches(3.20), Inches(5.20)]
    sup_w = Inches(6.13)
    sup_h = Inches(1.85)

    items = [
        ("Layer-ensemble WavLM distillation",
         "Three learnable softmax weights over WavLM layers {4, 8, 11} — "
         "beats single-layer choice (SKILL, Interspeech 2024). Phase-1 "
         "validated at the final-layer simplest form (+0.7 pp KWS clean)."),
        ("Multi-level ECAPA distillation",
         "Cosine on final embedding + MSE on block-2 and block-3 features "
         "(projected via 1×1 conv) — ~5 % relative EER reduction over "
         "final-only KD."),
        ("Phonetic confusable hard-negative mining",
         "TIMIT phone-confusion-weighted Levenshtein on G2P output, sampled "
         "at 3× during training — directly targets the FA < 1 / hr risk."),
        ("Considered & rejected alternatives",
         "Audio Mamba — 5× BC-ResNet params for equal clean acc., deploy "
         "immature.  ·  LEAF / SincNet — only PCEN actually trains; already "
         "covered.  ·  GRL adversarial — removes the speaker signal we need.  "
         "·  Deep cross-attn — 4× compute, breaks xRT budget."),
    ]
    for idx, (title, body) in enumerate(items):
        x = sup_x[idx % 2]
        y = sup_y[idx // 2]
        add_card(s, x, y, sup_w, sup_h)
        add_text(s, x + Inches(0.2), y + Inches(0.1), sup_w - Inches(0.4),
                 Inches(0.4),
                 [{"text": title, "size": 12, "bold": True, "color": SRI_DARK}])
        add_text(s, x + Inches(0.2), y + Inches(0.5), sup_w - Inches(0.4),
                 sup_h - Inches(0.6),
                 [{"text": body, "size": 11, "color": BODY_GRAY}])


def slide_10_datasets(prs, page):
    s = blank_slide(prs, "Open Datasets — Used & To Be Published", page)

    # Used table (top)
    add_text(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
             [label_para("To be used  (all permissive licences, downloadable)",
                          15)])
    add_card(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.7))

    headers = [("Dataset", Inches(3.4)),
               ("Role", Inches(6.6)),
               ("Licence", Inches(2.0))]
    rows = [
        ("Google Speech Commands v0.02",
         "KWS positives + general negatives", "CC-BY 4.0"),
        ("VoxCeleb1",
         "Speaker verification + impostor pool", "CC-BY 4.0"),
        ("LibriPhrase",
         "Joint phrase + speaker fine-tuning", "MIT"),
        ("MUSAN",
         "Crowd / babble / traffic noise injection (−5 to 30 dB)",
         "CC-BY 4.0"),
        ("OpenSLR RIRs",
         "Room impulse responses → 0.5–5 m distance simulation",
         "Apache-2.0"),
        ("MLCommons MSWC  (H1 eval)",
         "Multilingual zero-shot keyword generalisation (Tamil)",
         "CC-BY 4.0"),
    ]
    _render_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
                  headers, rows, row_h=Inches(0.45),
                  header_fill=SRI_BLUE)

    # Published card
    add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.4),
             [label_para("To be published  (Apache-2.0)", 15)])
    add_card(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.35),
             fill=ACCENT_BG_GOLD, border=SRI_GOLD)
    add_text(s, Inches(0.7), Inches(5.9), Inches(12.0), Inches(1.3),
             [{"runs": [
                 {"text": "Confusable-trigger benchmark.  ",
                  "size": 12, "bold": True, "color": SRI_GOLD},
                 {"text": "For any chosen trigger word, the mined top-50 "
                          "phonetic confusables synthesised in 100+ XTTS-v2 "
                          "voices + their RIR / MUSAN augmentations. Useful "
                          "for any future personalised-KWS work targeting "
                          "FA < 1 / hr. Released under Apache-2.0.",
                  "size": 11, "color": BODY_GRAY},
             ]}])


def slide_11_models(prs, page):
    s = blank_slide(prs, "Open Models — Used & Developed", page)

    # Used (teachers)
    add_text(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.4),
             [label_para("To be used  (frozen, distillation teachers)", 15)])
    add_card(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.6))
    headers = [("Model", Inches(3.0)),
               ("Role", Inches(6.0)),
               ("Source / Licence", Inches(3.0))]
    rows = [
        ("WavLM-Base+", "Phonetic / noise-robust teacher",
         "microsoft/wavlm-base-plus · MIT"),
        ("ECAPA-TDNN", "Speaker biometric teacher",
         "speechbrain/spkrec-ecapa · Apache-2.0"),
        ("XTTS-v2", "Multi-voice TTS for trigger augmentation",
         "coqui/XTTS-v2 · CPML"),
        ("g2p-en", "Trigger word → ARPAbet phones", "Apache-2.0"),
    ]
    _render_table(s, Inches(0.7), Inches(1.6), Inches(11.9),
                  headers, rows, row_h=Inches(0.45),
                  header_fill=SRI_BLUE)

    # Developed
    add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.4),
             [label_para("To be developed / trained from scratch  (Apache-2.0)",
                          15)])
    add_card(s, Inches(0.5), Inches(4.75), Inches(12.3), Inches(2.4))
    headers2 = [("Model", Inches(4.0)),
                ("Description", Inches(7.9))]
    rows2 = [
        ("BC-ResNet-8 (modified) joint student",
         "TRM branching, FiLM speaker conditioning, MQMHA SV pooling — "
         "joint KWS + SV, ~830 K params"),
        ("WavLM → KWS probing head",
         "2-layer MLP, ~50 K params, one-time training"),
        ("INT8-quantised ONNX export",
         "Edge-deployable artefact — 0.79 MB on disk"),
    ]
    _render_table(s, Inches(0.7), Inches(4.85), Inches(11.9),
                  headers2, rows2, row_h=Inches(0.65),
                  header_fill=SRI_GREEN)


def slide_12_tools(prs, page):
    s = blank_slide(prs, "AI / GenAI / Agentic Tools & Best Practices", page)

    # Tools (left column)
    add_text(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.4),
             [label_para("Tools used", 15)])
    add_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5))
    tools = [
        ("PyTorch 2.x",
         "+ torch.compile(mode=\"reduce-overhead\") for ~15 % Kaggle speedup."),
        ("HuggingFace Transformers / SpeechBrain",
         "frozen teacher inference (WavLM, ECAPA)."),
        ("Coqui XTTS-v2",
         "generative-AI component for synthetic speaker diversity in trigger "
         "and confusable data."),
        ("Optuna (TPE sampler)",
         "automated multi-task loss-weight + AAM hyperparameter search, "
         "scoped to Stage 2."),
        ("ONNX Runtime + Quantization Toolkit",
         "INT8 edge export and latency benchmarking."),
        ("Google kws_streaming",
         "reference for stateful streaming inference."),
        ("Claude Code (Anthropic CLI)",
         "Agentic AI helper and companion"),
    ]
    paras = []
    for name, desc in tools:
        paras.append({"runs": [
            {"text": f"• {name}  ", "size": 12, "bold": True,
             "color": SRI_BLUE},
            {"text": desc, "size": 11, "color": BODY_GRAY},
        ], "space_after": 4})
    add_text(s, Inches(0.7), Inches(1.6), Inches(5.7), Inches(5.35), paras)

    # Best practices (right column)
    add_text(s, Inches(6.8), Inches(1.05), Inches(6.0), Inches(0.4),
             [label_para("Best practices uncovered", 15)])
    add_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5))
    practices = [
        ("Symmetric augmentation in distillation.",
         "Teacher and student must see the same augmented audio, otherwise "
         "the student collapses to the teacher's clean manifold."),
        ("Multi-level KD beats single-layer KD",
         "for both phonetic (WavLM) and biometric (ECAPA) teachers — "
         "measured ~5 % relative metric gain."),
        ("Bake the deployment operating point into training",
         "(online τₛ calibration). Post-hoc thresholding cannot fix "
         "an embedding space that wasn't trained to be FAR-aware."),
        ("Structural conditioning > threshold gating",
         "for personalised triggers. FiLM-modulated KWS is materially safer "
         "for FA-bounded scenarios than \"AND of two independent classifiers\"."),
        ("Curriculum on SNR, not just on data difficulty.",
         "Starting at −5 dB diverges; ramp from 30 → −5 dB "
         "across the first 15 epochs."),
    ]
    paras2 = []
    for i, (head, tail) in enumerate(practices, 1):
        paras2.append({"runs": [
            {"text": f"{i}.  ", "size": 12, "bold": True, "color": SRI_RED},
            {"text": head + "  ", "size": 12, "bold": True, "color": SRI_DARK},
            {"text": tail, "size": 11, "color": BODY_GRAY},
        ], "space_after": 6})
    paras2.append({"runs": [
        {"text": "Licence:  ", "size": 11, "bold": True, "color": SRI_GREEN},
        {"text": "all code, weights, and the published benchmark released "
                 "under Apache-2.0, per hackathon rules.",
         "size": 11, "italic": True, "color": BODY_GRAY},
    ], "space_before": 4})
    add_text(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(5.35), paras2)


def slide_sv_branch_poc(prs, page):
    s = blank_slide(prs, "Phase-1 Empirical — SV Branch POC", page)

    # Left figure + caption
    add_image_fit(s, ILLUS_DIR / "8_sv_progression.png",
                  Inches(0.4), Inches(1.05), Inches(7.4), Inches(4.5))
    add_text(s, Inches(0.4), Inches(5.7), Inches(7.4), Inches(1.4),
             [{"text": "Trained on VoxCeleb1 dev (1 211 speakers, 24 K "
                       "utterances), evaluated on the canonical VoxCeleb1-O "
                       "trial list (37 720 same / different pairs). Student "
                       "is a 384 K-param BC-ResNet trunk + attentive-stat "
                       "pooling + 256-dim projection. Teacher: pretrained "
                       "resemblyzer (GE2E-LSTM).",
               "size": 10, "italic": True, "color": SRI_GRAY}])

    # Right: 3 stacked cards
    rx = Inches(8.0)
    rw = Inches(4.95)

    add_card(s, rx, Inches(1.05), rw, Inches(1.65),
             fill=ACCENT_BG_BLUE, border=SRI_BLUE)
    add_text(s, rx + Inches(0.2), Inches(1.1), rw - Inches(0.4), Inches(1.6),
             [{"text": "Current state", "size": 13, "bold": True,
               "color": SRI_BLUE, "space_after": 4},
              {"runs": [
                  {"text": "10.64 % EER", "size": 14, "bold": True,
                   "color": SRI_BLUE},
                  {"text": " on VoxCeleb1-O after 40 epochs joint "
                           "distillation + Sub-center AAM-Softmax. Student "
                           "is 384 K params — within budget for the joint "
                           "student.",
                   "size": 11, "color": BODY_GRAY},
              ]}])

    add_card(s, rx, Inches(2.85), rw, Inches(2.50),
             fill=ACCENT_BG_RED, border=SRI_RED)
    add_text(s, rx + Inches(0.2), Inches(2.9), rw - Inches(0.4), Inches(2.45),
             [{"text": "Why v1 collapsed at 28.67 %", "size": 13, "bold": True,
               "color": SRI_RED, "space_after": 4},
              {"text": "Pure cosine distillation with no classification "
                       "signal — embeddings cluster at a single point, no "
                       "angular discrimination.",
               "size": 11, "color": BODY_GRAY, "space_after": 4},
              {"runs": [
                  {"text": "Sub-center AAM with margin m = 0.2 ", "size": 11,
                   "color": BODY_GRAY},
                  {"text": "forces speakers into separate angular cones; ",
                   "size": 11, "color": BODY_GRAY},
                  {"text": "+18 pp gain", "size": 12, "bold": True,
                   "color": SRI_RED},
                  {"text": " from the loss change alone.", "size": 11,
                   "color": BODY_GRAY},
              ]}])

    add_card(s, rx, Inches(5.50), rw, Inches(1.55),
             fill=ACCENT_BG_GOLD, border=SRI_GOLD)
    add_text(s, rx + Inches(0.2), Inches(5.55), rw - Inches(0.4), Inches(1.5),
             [{"text": "Path to ≤ 5 % (v3)", "size": 13, "bold": True,
               "color": SRI_GOLD, "space_after": 4},
              {"text": "(i) Swap teacher to ECAPA-TDNN (1 % EER vs 4 %); "
                       "(ii) tune λ_distill down; (iii) train on full "
                       "VoxCeleb1 (~150 K utts) for 80 epochs.  Projected v3 "
                       "EER: 4–5 %  —  competitive with x-vector at ~10× "
                       "smaller.",
               "size": 11, "color": BODY_GRAY}])


def slide_film_rejection(prs, page):
    s = blank_slide(prs, "Phase-1 Empirical — FiLM Rejection Trade-off", page)

    # Left figure + caption
    add_image_fit(s, ILLUS_DIR / "9_film_tradeoff.png",
                  Inches(0.4), Inches(1.05), Inches(7.4), Inches(4.4))
    add_text(s, Inches(0.4), Inches(5.6), Inches(7.4), Inches(1.5),
             [{"text": "Joint KWS + SV + FiLM model (500 K params, 400 train "
                       "speakers, 8 K utts, 30 epochs). Eval on 28 held-out "
                       "speakers × 10 test items each, with 5-clip "
                       "enrollment averaging. Genuine = enrolled speaker "
                       "says keyword; impostor = different speaker enrolled, "
                       "keyword spoken by another. Collapse ratio = mean "
                       "P(true_kw)_impostor / mean P(true_kw)_genuine.",
               "size": 10, "italic": True, "color": SRI_GRAY}])

    rx = Inches(8.0)
    rw = Inches(4.95)

    add_card(s, rx, Inches(1.05), rw, Inches(2.05),
             fill=ACCENT_BG_BLUE, border=SRI_BLUE)
    add_text(s, rx + Inches(0.2), Inches(1.1), rw - Inches(0.4), Inches(2.0),
             [{"text": "Mechanism validated", "size": 13, "bold": True,
               "color": SRI_BLUE, "space_after": 4},
              {"runs": [
                  {"text": "95 % TAR + 86 % impostor-rejection", "size": 13,
                   "bold": True, "color": SRI_BLUE},
                  {"text": " at 14 % FAR when train + eval embedding "
                           "distributions are matched (v2 own-audio "
                           "diagnostic).",
                   "size": 11, "color": BODY_GRAY},
              ], "space_after": 4},
              {"runs": [
                  {"text": "Probability collapse ratio ", "size": 11,
                   "color": BODY_GRAY},
                  {"text": "0.20", "size": 13, "bold": True, "color": SRI_BLUE},
                  {"text": "  —  impostors get ~5× less keyword confidence "
                           "than genuine.",
                   "size": 11, "color": BODY_GRAY},
              ]}])

    add_card(s, rx, Inches(3.25), rw, Inches(2.10),
             fill=ACCENT_BG_RED, border=SRI_RED)
    add_text(s, rx + Inches(0.2), Inches(3.3), rw - Inches(0.4), Inches(2.05),
             [{"text": "Trade-off characterised", "size": 13, "bold": True,
               "color": SRI_RED, "space_after": 4},
              {"runs": [
                  {"text": "Closing the train / eval gap (v3 paired "
                           "enrollment) lifts TAR to 92 % but loosens "
                           "rejection — FAR jumps to ",
                   "size": 11, "color": BODY_GRAY},
                  {"text": "48 %", "size": 13, "bold": True, "color": SRI_RED},
                  {"text": ".  The system has a tunable TAR / FAR knob via "
                           "impostor-loss weighting and hard-negative mining.",
                   "size": 11, "color": BODY_GRAY},
              ]}])

    add_card(s, rx, Inches(5.50), rw, Inches(1.55),
             fill=ACCENT_BG_GOLD, border=SRI_GOLD)
    add_text(s, rx + Inches(0.2), Inches(5.55), rw - Inches(0.4), Inches(1.5),
             [{"text": "Path to deployment (v4)", "size": 13, "bold": True,
               "color": SRI_GOLD, "space_after": 4},
              {"text": "Hard-negative mining + λ_imp = 2 + max-impostor-ratio "
                       "0.7.  Target: ≥ 85 % TAR + ≤ 25 % FAR + ≥ 70 % REJ "
                       "at the production operating point.",
               "size": 11, "color": BODY_GRAY}])


def slide_far_hero(prs, page):
    s = blank_slide(prs, "Novelty & Innovation — FAR-in-the-Loop Training "
                          "(hero)", page)

    # Wide figure banner at top
    add_image_fit(s, ILLUS_DIR / "3_far_in_loop.png",
                  Inches(0.4), Inches(1.05), Inches(12.5), Inches(2.4))

    # 3 column blocks below
    col_y = Inches(3.65)
    col_h = Inches(3.3)
    col_w = Inches(4.0)
    col_gap = Inches(0.13)
    col_x = [Inches(0.5),
             Inches(0.5) + col_w + col_gap,
             Inches(0.5) + (col_w + col_gap) * 2]

    blocks = [
        ("Core invention", SRI_RED, ACCENT_BG_RED, [
            "The deployment false-accept threshold τₛ is re-estimated every "
            "500 training steps on a held-out impostor pool.",
            "",
            "The resulting FRR-at-target-FAR is fed back as a differentiable "
            "soft penalty in the joint loss.",
        ]),
        ("Why it matters", SRI_BLUE, ACCENT_BG_BLUE, [
            "Conventional pipelines treat τₛ as a post-hoc deployment knob.",
            "",
            "If the training loss never sees the FA < 1 / hr constraint, the "
            "embedding geometry cannot cluster tightly enough for any "
            "threshold to satisfy it.",
            "",
            "Ref: PK-MTL — vanilla BC-ResNet hits 64.32 % FRR in target-only "
            "KWS; task-adaptive training brings this to 6.56 %.",
        ]),
        ("Result", SRI_GREEN, ACCENT_BG_GREEN, [
            "The operating point is baked into the embedding geometry — "
            "not tuned around it.",
            "",
            "FA < 1 / hr becomes a property of the trained model, not of a "
            "post-hoc threshold sweep.",
        ]),
    ]
    for x, (title, color, bg, paras) in zip(col_x, blocks):
        add_card(s, x, col_y, col_w, col_h, fill=bg, border=color)
        add_text(s, x + Inches(0.2), col_y + Inches(0.12),
                 col_w - Inches(0.4), Inches(0.45),
                 [{"text": title, "size": 13, "bold": True, "color": color}])
        body = []
        for line in paras:
            if line == "":
                body.append({"text": "", "size": 6, "color": BODY_GRAY,
                             "space_after": 2})
            else:
                size = 10 if line.startswith("Ref:") else 11
                italic = line.startswith("Ref:")
                body.append({"text": line, "size": size, "italic": italic,
                             "color": BODY_GRAY, "space_after": 4})
        add_text(s, x + Inches(0.25), col_y + Inches(0.6),
                 col_w - Inches(0.5), col_h - Inches(0.75), body)


def slide_demo_plan(prs, page):
    s = blank_slide(prs, "Deployment & Demo Plan (for Grand Finale)", page)

    # Left: enrollment runtime figure (portrait — ar 0.72)
    add_image_fit(s, ILLUS_DIR / "5_enrollment_runtime.png",
                  Inches(0.4), Inches(1.05), Inches(5.2), Inches(5.9))
    add_text(s, Inches(0.4), Inches(7.0), Inches(5.2), Inches(0.4),
             [{"text": "Enrollment → runtime flow (5 utts ~30 s → "
                       "stored ID prototype → online inference).",
               "size": 9, "italic": True, "color": SRI_GRAY, "align": "center"}])

    # Right: 2 cards (demo script + hardware)
    rx = Inches(5.9)
    rw = Inches(7.0)

    add_card(s, rx, Inches(1.05), rw, Inches(4.6),
             fill=ACCENT_BG_BLUE, border=SRI_BLUE)
    add_text(s, rx + Inches(0.25), Inches(1.15), rw - Inches(0.5), Inches(4.5),
             [{"text": "Live demo script", "size": 14, "bold": True,
               "color": SRI_BLUE, "space_after": 8},
              {"runs": [
                  {"text": "1.  ", "size": 12, "bold": True, "color": SRI_BLUE},
                  {"text": "Enroll a live judge:  ", "size": 12, "bold": True,
                   "color": SRI_DARK},
                  {"text": "5 clean utterances (~30 s) → stored speaker "
                           "prototype.", "size": 12, "color": BODY_GRAY},
              ], "space_after": 8},
              {"runs": [
                  {"text": "2.  ", "size": 12, "bold": True, "color": SRI_BLUE},
                  {"text": "Clean trigger ", "size": 12, "bold": True,
                   "color": SRI_DARK},
                  {"text": "from the judge   →   ", "size": 12,
                   "color": BODY_GRAY},
                  {"text": "fires.", "size": 12, "bold": True,
                   "color": SRI_GREEN},
              ], "space_after": 8},
              {"runs": [
                  {"text": "3.  ", "size": 12, "bold": True, "color": SRI_BLUE},
                  {"text": "Noisy trigger ", "size": 12, "bold": True,
                   "color": SRI_DARK},
                  {"text": "— traffic at −5 dB from phone speaker   →   ",
                   "size": 12, "color": BODY_GRAY},
                  {"text": "fires.", "size": 12, "bold": True,
                   "color": SRI_GREEN},
              ], "space_after": 8},
              {"runs": [
                  {"text": "4.  ", "size": 12, "bold": True, "color": SRI_BLUE},
                  {"text": "Impostor ", "size": 12, "bold": True,
                   "color": SRI_DARK},
                  {"text": "(teammate) says the same word   →   ",
                   "size": 12, "color": BODY_GRAY},
                  {"text": "rejected.", "size": 12, "bold": True,
                   "color": SRI_RED},
              ], "space_after": 8},
              {"runs": [
                  {"text": "5.  ", "size": 12, "bold": True, "color": SRI_BLUE},
                  {"text": "Confusable ", "size": 12, "bold": True,
                   "color": SRI_DARK},
                  {"text": "phrase (\"turn on\" vs \"turn off\")   →   ",
                   "size": 12, "color": BODY_GRAY},
                  {"text": "rejected.", "size": 12, "bold": True,
                   "color": SRI_RED},
              ]}])

    add_card(s, rx, Inches(5.85), rw, Inches(1.45),
             fill=ACCENT_BG_GOLD, border=SRI_GOLD)
    add_text(s, rx + Inches(0.25), Inches(5.95), rw - Inches(0.5), Inches(1.35),
             [{"text": "Hardware", "size": 13, "bold": True,
               "color": SRI_GOLD, "space_after": 4},
              {"runs": [
                  {"text": "Raspberry Pi 4B ", "size": 12, "bold": True,
                   "color": SRI_DARK},
                  {"text": "(conservative spec) ", "size": 11,
                   "color": BODY_GRAY},
                  {"text": "or ", "size": 12, "color": BODY_GRAY},
                  {"text": "Galaxy-class NPU via QNN", "size": 12, "bold": True,
                   "color": SRI_DARK},
                  {"text": ".  No network connection during demo — entirely "
                           "on-device.",
                   "size": 11, "color": BODY_GRAY},
              ]}])


def slide_closing(prs, page):
    """Standout closing slide — navy background, white text."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, SRI_DARK)

    # Gold accent bar
    add_rect(s, 0, Inches(2.6), SLIDE_W, Inches(0.08), SRI_GOLD)

    # Big headline
    add_text(s, Inches(0.8), Inches(1.0), SLIDE_W - Inches(1.6), Inches(1.5),
             [{"text": "Closing the FA < 1 / hr gap on < 3 M params",
               "size": 36, "bold": True,
               "color": RGBColor(0xFF, 0xFF, 0xFF), "align": "center"}],
             anchor="middle")

    # Tagline pillars
    add_text(s, Inches(0.8), Inches(3.0), SLIDE_W - Inches(1.6), Inches(0.6),
             [{"runs": [
                 {"text": "FAR-in-the-Loop training", "size": 18,
                  "bold": True, "color": RGBColor(0xFF, 0xFF, 0xFF)},
                 {"text": "    ·    ", "size": 18,
                  "color": SRI_GOLD},
                 {"text": "FiLM-conditioned KWS", "size": 18,
                  "bold": True, "color": RGBColor(0xFF, 0xFF, 0xFF)},
                 {"text": "    ·    ", "size": 18,
                  "color": SRI_GOLD},
                 {"text": "Multi-teacher distillation", "size": 18,
                  "bold": True, "color": RGBColor(0xFF, 0xFF, 0xFF)},
             ], "align": "center"}])

    # Measured anchors
    add_text(s, Inches(0.8), Inches(4.0), SLIDE_W - Inches(1.6), Inches(1.6),
             [{"runs": [
                 {"text": "Phase-1 measured  ·  ", "size": 14, "italic": True,
                  "color": RGBColor(0xCC, 0xD3, 0xE0)},
                 {"text": "KWS clean 98.3 %", "size": 16, "bold": True,
                  "color": SRI_GOLD},
                 {"text": "  ·  ", "size": 14,
                  "color": RGBColor(0xCC, 0xD3, 0xE0)},
                 {"text": "English EER 12.4 %", "size": 16, "bold": True,
                  "color": SRI_GOLD},
                 {"text": "  ·  ", "size": 14,
                  "color": RGBColor(0xCC, 0xD3, 0xE0)},
                 {"text": "Tamil zero-shot 14.7 %", "size": 16, "bold": True,
                  "color": SRI_GOLD},
                 {"text": "  ·  ", "size": 14,
                  "color": RGBColor(0xCC, 0xD3, 0xE0)},
                 {"text": "Pi 4B INT8 12 ms", "size": 16, "bold": True,
                  "color": SRI_GOLD},
                 {"text": "  ·  ", "size": 14,
                  "color": RGBColor(0xCC, 0xD3, 0xE0)},
                 {"text": "499 K params (17 % cap)", "size": 16, "bold": True,
                  "color": SRI_GOLD},
             ], "align": "center"}])

    # Bottom credit line
    add_text(s, Inches(0.8), Inches(6.6), SLIDE_W - Inches(1.6), Inches(0.5),
             [{"text": "Team MalayM09  ·  Malay Mishra  ·  BITS Pilani  ·  "
                       "Problem PS04  ·  Apache-2.0",
               "size": 13, "italic": True,
               "color": RGBColor(0xCC, 0xD3, 0xE0), "align": "center"}])


def slide_glossary(prs, page):
    """Appendix — 6-card glossary grouped by category."""
    s = blank_slide(prs, "Appendix — Glossary", page)

    # 3 columns × 2 rows of cards
    card_w = Inches(4.10)
    card_h = Inches(2.95)
    gap_x = Inches(0.13)
    gap_y = Inches(0.13)
    origin_x = Inches(0.45)
    origin_y = Inches(1.05)
    positions = [
        (origin_x + (card_w + gap_x) * c, origin_y + (card_h + gap_y) * r)
        for r in (0, 1) for c in (0, 1, 2)
    ]

    categories = [
        ("Frontend & Architecture", SRI_BLUE, ACCENT_BG_BLUE, [
            ("PCEN",      "Per-Channel Energy Normalization — learnable frontend"),
            ("LEAF",      "Learnable Auditory Frontend"),
            ("BC-ResNet", "Broadcast Residual Network — compact CNN for KWS"),
            ("TRM",       "Task Representation Module — branch decoupler"),
            ("MQMHA",     "Multi-Query Multi-Head Attention pooling"),
            ("FiLM",      "Feature-wise Linear Modulation (γ scale + β shift)"),
            ("stop-grad", "edge that blocks back-propagation"),
        ]),
        ("Losses & Training", SRI_RED, ACCENT_BG_RED, [
            ("KD",            "Knowledge Distillation"),
            ("Focal CE",      "Focal Cross-Entropy"),
            ("AAM-Softmax",   "Additive Angular Margin Softmax"),
            ("Sub-center",    "AAM variant — handles label noise (K=3 cones)"),
            ("KL",            "Kullback–Leibler divergence"),
            ("L_hint",        "feature-matching distillation loss"),
            ("T,  λᵢ,  m",    "distill temperature  ·  loss weights  ·  margin"),
            ("GRL",           "Gradient Reversal Layer"),
        ]),
        ("Metrics", SRI_GREEN, ACCENT_BG_GREEN, [
            ("KWS, SV",     "Keyword Spotting  ·  Speaker Verification"),
            ("TA / FA",     "True / False Acceptance (event)"),
            ("TAR / FAR",   "True / False Acceptance Rate (%)"),
            ("FRR",         "False Rejection Rate"),
            ("EER",         "Equal Error Rate (FAR = FRR crossover)"),
            ("AUC",         "Area Under ROC Curve"),
            ("REJ",         "impostor-rejection rate"),
            ("pp,  xRT,  SNR", "percentage pts  ·  real-time factor  ·  signal/noise (dB)"),
        ]),
        ("Datasets & Models", SRI_GOLD, ACCENT_BG_GOLD, [
            ("SC V2",        "Google Speech Commands v0.02 (KWS corpus)"),
            ("MSWC",         "MLCommons Spoken Words Corpus (multilingual)"),
            ("VoxCeleb1-O",  "Original VoxCeleb1 trial list (37 720 pairs)"),
            ("MUSAN",        "Music-Speech-Noise corpus (noise augmentation)"),
            ("RIR",          "Room Impulse Response (far-field simulation)"),
            ("WavLM",        "Microsoft pretrained speech encoder"),
            ("ECAPA-TDNN",   "pretrained speaker-verification model"),
            ("XTTS-v2 / g2p-en", "Coqui multilingual TTS  ·  English G2P"),
        ]),
        ("Deployment", SRI_DARK, RGBColor(0xEE, 0xF1, 0xF7), [
            ("ONNX",  "Open Neural Network Exchange (export format)"),
            ("ORT",   "ONNX Runtime — inference engine"),
            ("PTQ",   "Post-Training Quantization"),
            ("QAT",   "Quantization-Aware Training"),
            ("INT8",  "8-bit integer (vs FP32)"),
            ("Pi 4B", "Raspberry Pi 4 Model B — conservative edge target"),
            ("QNN",   "Qualcomm Neural Network SDK"),
            ("NPU",   "Neural Processing Unit (on-device accelerator)"),
        ]),
        ("Cited Work & Misc", SRI_BLUE, ACCENT_BG_BLUE, [
            ("SKILL",     "layer-ensemble distillation, Interspeech 2024"),
            ("PK-MTL",    "Per-Keyword Multi-Task Learning  (arXiv:2206.13708)"),
            ("GE2E",      "Generalized End-to-End speaker loss"),
            ("TPE",       "Tree-structured Parzen Estimator (Optuna sampler)"),
            ("LLM",       "Large Language Model"),
            ("τₛ",        "speaker decision threshold (FAR-in-the-Loop target)"),
            ("Apache-2.0, CC-BY, MIT, CPML", "open-source licence shorthands"),
        ]),
    ]

    for (x, y), (title, color, bg, entries) in zip(positions, categories):
        add_card(s, x, y, card_w, card_h, fill=bg, border=color)
        # Category title
        add_text(s, x + Inches(0.18), y + Inches(0.08),
                 card_w - Inches(0.36), Inches(0.32),
                 [{"text": title, "size": 12, "bold": True, "color": color}])
        # Entries
        paras = []
        for term, defn in entries:
            paras.append({"runs": [
                {"text": term, "size": 9.5, "bold": True, "color": SRI_DARK},
                {"text": "   " + defn, "size": 9, "color": BODY_GRAY},
            ], "space_after": 2})
        add_text(s, x + Inches(0.20), y + Inches(0.44),
                 card_w - Inches(0.4), card_h - Inches(0.55),
                 paras, margin=0.04)

    # Bottom note
    add_text(s, Inches(0.45), Inches(7.18),
             SLIDE_W - Inches(0.9), Inches(0.28),
             [{"text": "Reference appendix — terms used across Slides 2–17. "
                       "Bolded acronym  ·  full expansion  ·  context.",
               "size": 9, "italic": True, "color": SRI_GRAY,
               "align": "center"}])


# ──────────────────────────────────────────────────────────────────────────────
# TABLE HELPER
# ──────────────────────────────────────────────────────────────────────────────
def _render_table(slide, x, y, total_w, headers, rows,
                  row_h, header_fill=SRI_BLUE):
    # Header row
    cx = x
    add_rect(slide, x, y, total_w, Inches(0.42), header_fill)
    for label, col_w in headers:
        add_text(slide, cx + Inches(0.1), y, col_w - Inches(0.1),
                 Inches(0.42),
                 [{"text": label, "size": 11, "bold": True,
                   "color": RGBColor(0xFF, 0xFF, 0xFF)}], anchor="middle")
        cx += col_w
    ry = y + Inches(0.45)
    for i, row in enumerate(rows):
        if i % 2 == 0:
            add_rect(slide, x, ry, total_w, row_h,
                     RGBColor(0xF4, 0xF7, 0xFB))
        cx = x
        for (label, col_w), val in zip(headers, row):
            add_text(slide, cx + Inches(0.1), ry, col_w - Inches(0.1),
                     row_h,
                     [{"text": val, "size": 11,
                       "color": SRI_DARK if label == headers[0][0] else BODY_GRAY,
                       "bold": label == headers[0][0]}],
                     anchor="middle")
            cx += col_w
        ry += row_h


# ──────────────────────────────────────────────────────────────────────────────
# BUILD
# ──────────────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_team(prs, 1)
    slide_2_problem(prs, 2)
    slide_3_solution_high(prs, 3)
    slide_4_architecture(prs, 4)
    slide_5_distillation(prs, 5)
    slide_6_params_stages(prs, 6)
    slide_7_phase1_kws(prs, 7)
    slide_sv_branch_poc(prs, 8)
    slide_film_rejection(prs, 9)
    slide_8_phase1_deploy_babble(prs, 10)
    slide_far_hero(prs, 11)
    slide_9_novelty(prs, 12)
    slide_10_datasets(prs, 13)
    slide_11_models(prs, 14)
    slide_demo_plan(prs, 15)
    slide_12_tools(prs, 16)
    slide_closing(prs, 17)
    slide_glossary(prs, 18)

    out = Path(__file__).resolve().parent / "AX_Hackathon_Phase1_Blueprint_FILLED.pptx"
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    build()
